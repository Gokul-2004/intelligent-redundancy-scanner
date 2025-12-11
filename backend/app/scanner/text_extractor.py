"""Extract text from various file types for content-based duplicate detection"""
import io
from typing import Optional
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def extract_text_from_file(content: bytes, mime_type: str, filename: str) -> Optional[str]:
    """
    Extract text content from file based on MIME type
    
    Args:
        content: File content as bytes
        mime_type: MIME type of the file
        filename: Filename (for fallback detection)
    
    Returns:
        Extracted text or None if extraction fails
    """
    try:
        # PDF files
        if mime_type == 'application/pdf' or filename.lower().endswith('.pdf'):
            return extract_text_from_pdf(content)
        
        # Word documents
        elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                          'application/msword'] or filename.lower().endswith(('.docx', '.doc')):
            return extract_text_from_docx(content)
        
        # Excel files
        elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                          'application/vnd.ms-excel'] or filename.lower().endswith(('.xlsx', '.xls')):
            return extract_text_from_xlsx(content)
        
        # PowerPoint
        elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                          'application/vnd.ms-powerpoint'] or filename.lower().endswith(('.pptx', '.ppt')):
            return extract_text_from_pptx(content)
        
        # Plain text
        elif mime_type == 'text/plain' or filename.lower().endswith('.txt'):
            try:
                return content.decode('utf-8')
            except:
                try:
                    return content.decode('latin-1')
                except:
                    return None
        
        # HTML
        elif mime_type == 'text/html' or filename.lower().endswith(('.html', '.htm')):
            # Basic HTML text extraction (remove tags)
            try:
                text = content.decode('utf-8')
                # Simple tag removal (basic implementation)
                import re
                text = re.sub(r'<[^>]+>', '', text)
                return text.strip()
            except:
                return None
        
        return None
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        return None


def extract_text_from_pdf(content: bytes) -> Optional[str]:
    """Extract text from PDF"""
    try:
        pdf_file = io.BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text_parts = []
        
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        return '\n'.join(text_parts) if text_parts else None
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return None


def extract_text_from_docx(content: bytes) -> Optional[str]:
    """Extract text from DOCX"""
    try:
        doc_file = io.BytesIO(content)
        doc = Document(doc_file)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' '.join(cell.text for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts) if text_parts else None
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        return None


def extract_text_from_xlsx(content: bytes) -> Optional[str]:
    """Extract text from XLSX"""
    try:
        xlsx_file = io.BytesIO(content)
        workbook = load_workbook(xlsx_file, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = []
            
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    sheet_text.append(row_text)
            
            if sheet_text:
                text_parts.append(f"Sheet: {sheet_name}\n" + '\n'.join(sheet_text))
        
        return '\n\n'.join(text_parts) if text_parts else None
    except Exception as e:
        print(f"XLSX extraction error: {e}")
        return None


def extract_text_from_pptx(content: bytes) -> Optional[str]:
    """Extract text from PPTX"""
    try:
        pptx_file = io.BytesIO(content)
        presentation = Presentation(pptx_file)
        text_parts = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
        
        return '\n\n'.join(text_parts) if text_parts else None
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return None


def normalize_text(text: str) -> str:
    """Normalize text for comparison (lowercase, remove extra whitespace)"""
    if not text:
        return ""
    # Convert to lowercase, normalize whitespace
    import re
    text = text.lower()
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

